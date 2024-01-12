import pptx
import time
import pandas as pd
import comtypes.client
import os
import email
import smtplib
import ssl
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

# Email credentials and file names
sender_address = 'user@gmail.com'
sender_pass = 'pass'
presentation_name = 'certificate.pptx'
excel_file_name = 'Participant_Data.xlsx'
sheet_name = 'Sheet1'

def send_email(send_ID, pdf_name, person_name, pdf_suffix, fname):
    # Create an email message
    message = MIMEMultipart()
    message['From'] = sender_address
    message['To'] = send_ID
    message['Subject'] = 'Certificate for Participating in Competition'
    
    # Email content
    mail_content = f"""
    Dear {person_name},

    Thank you for participating in the competition.

    Find attached your certificate.

    Wishing you the best for all future endeavors	

    Regards,

    {fname}
    Coordinator - competition	    
    """
    message.attach(MIMEText(mail_content, "plain"))

    # Attach the certificate PDF
    with open(pdf_name, "rb") as attachment:
        part = MIMEBase("application", "octet-stream")
        part.set_payload(attachment.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", f"attachment; filename={pdf_suffix}")
        message.attach(part)

    # Attach an additional file (OngoingOnlineSessions.jpeg)
    with open("CM_Attachment.jpeg", "rb") as attachment:
        part = MIMEBase("application", "octet-stream")
        part.set_payload(attachment.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", f"attachment; filename=OngoingOnlineSessions.jpeg")
        message.attach(part)

    # Convert the message to a string
    text = message.as_string()

    # Log in to the server using secure context and send the email
    context = ssl.create_default_context()
    with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as server:
        server.login(sender_address, sender_pass)
        server.sendmail(sender_address, send_ID, text)
    print(f'Mail Sent to {person_name}')

def PPTtoPDF(inputFileName, outputFileName, formatType=32):
    # Convert PowerPoint presentation to PDF
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType)  # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

# Load PowerPoint presentation
prs = pptx.Presentation(presentation_name)
slide_layout = prs.slide_layouts[1]
slide = prs.slides[0]
participant_name_id = 0

# Load participant data from Excel
df = pd.read_excel(excel_file_name, sheet_name=sheet_name, engine='openpyxl')
col_names = df.columns
ppt_element_text = [f"<{i}>" for i in col_names]
ppt_element_id = [-1 for i in col_names]

# Identify the placeholders in the PowerPoint slide
for i in range(len(slide.shapes)):
    shape = slide.shapes[i]
    if not shape.has_text_frame:
        continue
    elif(shape.text in ppt_element_text):
        index_ = ppt_element_text.index(shape.text)
        ppt_element_id[index_] = i

# Replace placeholders with actual data for each participant
for person in df.index:
    for i, index_ in enumerate(ppt_element_id):
        if index_ == -1 and df.columns[i] != "Email":
            print(f"There is no field in the PowerPoint called {df.columns[i]}")
            continue
        paragraph = slide.shapes[index_].text_frame.paragraphs[0]
        p = paragraph._p
        for idx, run in enumerate(paragraph.runs):
            if idx == 0:
                continue
            p.remove(run._r)
        paragraph.runs[0].text = df[col_names[i]][person]
    pptx_name = os.getcwd() + '/Certificate_' + "_".join(str(df["Name"][person]).split(' ')) + '.pptx'
    pdf_name = os.getcwd() + '/Certificate_' + "_".join(str(df["Name"][person]).split(' ')) + '.pdf'
    prs.save(pptx_name)
    time.sleep(0.5)
    PPTtoPDF(pptx_name, pdf_name)
    pdf_name_suffix = pdf_name.split('/')[1]
    time.sleep(1)
    os.remove(pptx_name)
    # Uncomment the line below to send emails
    # send_email(df["Email"][person], pdf_name, df["Name"][person], pdf_name_suffix, df["fname"][person])
